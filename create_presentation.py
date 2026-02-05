from pptx import Presentationdef create_presentation():
    prs = Presentation()

    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = "Studyspace"
    if 1 in slide.placeholders:
    title.text = "Studyspace"
    if 1 in slide.placeholders:
        subtitle = slide.placeholders[1]
        subtitle.text = "Team Name: [Insert Team Name]\nTrack: EdTech (Smart Campus)"    def add_content_slide(heading, content_text):
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = heading
        
        tf = body_shape.text_frame
        tf.word_wrap = True
        
        # Split content by lines and add as bullets
        lines = content_text.strip().split('\n')
        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
                p.text = line.strip('- ').strip()
                p.level = 0
            else:
                p = tf.add_paragraph()
                p.text = line.strip('- ').strip()
                p.level = 0
    # Slide 2: Team Members
    team_content = """Harshit Siraswal - [Roll No] - [Branch & Year]
[Member 2 Name] - [Roll No] - [Branch & Year]
[Member 3 Name] - [Roll No] - [Branch & Year]
[Member 4 Name] - [Roll No] - [Branch & Year]"""
    add_content_slide("Team Members' Details", team_content)

    # Slide 3: Project Overview (Problem)
    problem_content = """Problem: Chaos in Academic Communication
- Fragmentation: Resources scattered across WhatsApp, Telegram, and Drive.
- Lack of Verification: Anonymous users lead to spam and misinformation.
- Ephemeral Nature: Valuable discussions are lost in chat streams.
- Resource Clutter: Difficulty finding specific notes or syllabus during exams."""
    add_content_slide("Project Overview: Problem Statement", problem_content)

    # Slide 4: Project Overview (Solution)
    solution_content = """Solution: A Verified, Structured Academic Ecosystem
- Identity Verification: Access strictly tied to college domains (@kiet.edu).
- Structured Chat Rooms: Persistent, topic-focussed discussions.
- Centralized Resources: Searchable repository for notes and syllabus.
- Sustainable Model: Freemium access for power users and clubs."""
    add_content_slide("Project Overview: Proposed Solution", solution_content)

    # Slide 5: Workflow / Architecture
    workflow_content = """System Architecture:
- Frontend: iOS Native App (Swift / SwiftUI)
- Authentication: Secure Login via Supabase Auth (College Email)
- Backend: Supabase (PostgreSQL Database + Realtime subscriptions)
- Storage: Centralized Document Storage for Resources
(Please insert flow diagrams/architecture figures here)"""
    add_content_slide("System Architecture", workflow_content)

    # Slide 6: Technology Stack
    tech_content = """Must Include:
- Swift & iOS App Development (Xcode, iOS SDK)
- SwiftUI for modern UI/UX

Additional Resources:
- Supabase (Backend as a Service)
- PostgreSQL (Database)
- Razorpay (Payment Gateway Integration)
- AI Resources (if applicable)"""
    add_content_slide("Technology Stack", tech_content)

    # Slide 7: Expected Impact
    impact_content = """Real-World Relevance:
- Unified Campus Communication: Brings students and faculty to one platform.
- Academic Integrity: Verified users ensure high-quality interactions.
- Enhanced Productivity: Study timers and organized rooms help focus.
- Beneficiaries: Students, Professors, Student Clubs, and Administration."""
    add_content_slide("Expected Impact / Use Case", impact_content)

    # End Slide
    end_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(end_slide_layout)
    title = slide.shapes.title
    title.text = "Thank You"
    subtitle = slide.placeholders[1]
    subtitle.text = "Q & A"
    # Save
    output_filename = "TeamName_EdTech.pptx"
    try:
        prs.save(output_filename)
        print(f"Presentation created successfully: {output_filename}")
    except PermissionError:
        print(f"Error: Cannot save '{output_filename}'. File may be open or directory not writable.")
    try:
        prs.save(output_filename)
        print(f"Presentation created successfully: {output_filename}")
    except PermissionError:
        print(f"Error: Cannot save '{output_filename}'. File may be open or directory not writable.")
    except Exception as e:
        print(f"Error saving presentation: {e}")